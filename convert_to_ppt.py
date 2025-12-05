#!/usr/bin/env python3
"""
Convert Reveal.js HTML presentation to PowerPoint (PPTX)
Improved version with images and better alignment
"""

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import re
from PIL import Image

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    if not hex_color or hex_color == 'none':
        return (255, 255, 255)  # Default to white
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3:
        hex_color = ''.join([c*2 for c in hex_color])
    try:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    except:
        return (255, 255, 255)

def clean_text(text):
    """Clean and format text"""
    if not text:
        return ""
    # Remove extra whitespace
    text = ' '.join(text.split())
    # Remove emoji-like characters that might not render well
    text = re.sub(r'[🏆⚠️🔍💡🌐📊]', '', text)
    return text.strip()

def find_images_in_section(section_soup, images_dir='images'):
    """Find all image references in a section"""
    images = []
    img_tags = section_soup.find_all('img')
    
    for img in img_tags:
        src = img.get('src', '')
        if src:
            # Handle different path formats
            img_path = None
            
            # Try different path combinations
            if src.startswith('images/'):
                img_path = os.path.join(images_dir, src.replace('images/', ''))
            elif src.startswith('/images/'):
                img_path = os.path.join(images_dir, src.replace('/images/', ''))
            elif os.path.isabs(src):
                img_path = src
            else:
                # Try relative to images directory
                img_path = os.path.join(images_dir, src)
            
            # Also try with just the filename
            if not img_path or not os.path.exists(img_path):
                filename = os.path.basename(src)
                img_path = os.path.join(images_dir, filename)
            
            # Check if file exists
            if img_path and os.path.exists(img_path):
                images.append({
                    'path': os.path.abspath(img_path),  # Use absolute path
                    'alt': img.get('alt', ''),
                    'style': img.get('style', '')
                })
                print(f"      Found image: {img_path}")
            else:
                print(f"      Warning: Image not found: {src} (tried: {img_path})")
    
    return images

def get_image_size(img_path, max_width=Inches(8), max_height=Inches(5)):
    """Get image dimensions, scaled to fit within max dimensions"""
    try:
        with Image.open(img_path) as img:
            # Get pixel dimensions
            pixel_width, pixel_height = img.size
            # Calculate aspect ratio
            aspect = pixel_width / pixel_height
            
            # Convert max dimensions from Inches to inches (float)
            if hasattr(max_width, 'inches'):
                max_w_inches = max_width.inches
            else:
                max_w_inches = float(max_width) / 914400 if max_width > 1000 else float(max_width)
            
            if hasattr(max_height, 'inches'):
                max_h_inches = max_height.inches
            else:
                max_h_inches = float(max_height) / 914400 if max_height > 1000 else float(max_height)
            
            # Convert pixels to inches (assuming 96 DPI standard)
            dpi = img.info.get('dpi', (96, 96))[0] if 'dpi' in img.info else 96
            width_inches = pixel_width / dpi
            height_inches = pixel_height / dpi
            
            # Scale to fit within max dimensions
            if width_inches > max_w_inches or height_inches > max_h_inches:
                if aspect > 1:  # Landscape
                    width_inches = max_w_inches
                    height_inches = max_w_inches / aspect
                else:  # Portrait
                    height_inches = max_h_inches
                    width_inches = max_h_inches * aspect
            
            # Return as Inches objects
            return Inches(width_inches), Inches(height_inches)
    except Exception as e:
        print(f"      Warning: Could not read image {img_path}: {e}")
        # Default size if image can't be read
        return Inches(6), Inches(4)

def parse_html_slides(html_file):
    """Parse HTML file and extract slide content"""
    with open(html_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    soup = BeautifulSoup(content, 'html.parser')
    slides = []
    
    # Find all section elements (slides)
    sections = soup.find_all('section')
    
    for section in sections:
        slide_data = {
            'title': '',
            'content': [],
            'background_color': None,
            'html_content': str(section),
            'soup': section,
            'images': []
        }
        
        # Get background color if specified
        if section.get('data-background-color'):
            slide_data['background_color'] = section.get('data-background-color')
        
        # Extract title (h1 or h2)
        title_elem = section.find(['h1', 'h2'])
        if title_elem:
            slide_data['title'] = clean_text(title_elem.get_text())
        
        # Find images
        slide_data['images'] = find_images_in_section(section)
        
        slides.append(slide_data)
    
    return slides

def add_text_element(slide, text, left, top, width, height, font_size=12, bold=False, 
                    color=None, alignment=PP_ALIGN.LEFT, wrap=True):
    """Add text to a slide at specific position"""
    if not text or not text.strip():
        return None
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = wrap
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    p = text_frame.paragraphs[0]
    p.text = clean_text(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.alignment = alignment
    
    if color:
        try:
            r, g, b = hex_to_rgb(color)
            p.font.color.rgb = RGBColor(r, g, b)
        except:
            pass
    
    return textbox

def process_slide_content(slide, slide_data):
    """Process and add content from HTML to slide"""
    soup = slide_data['soup']
    y_pos = Inches(1.2)
    max_y = Inches(6.8)
    left_margin = Inches(0.5)
    content_width = Inches(9)
    
    # Check for special slide types
    is_title_slide = slide_data['title'] and 'Comparative Analysis' in slide_data['title']
    is_hook_slide = '537 Million' in str(soup) or (slide_data['title'] and '537' in slide_data['title'])
    is_mission_slide = slide_data['title'] and 'Mission' in slide_data['title']
    is_data_slide = slide_data['title'] and ('Our Data' in slide_data['title'] or 'Pima Indians' in slide_data['title'])
    is_class_imbalance_slide = slide_data['title'] and ('Class Imbalance' in slide_data['title'] or 'Red Flag' in slide_data['title'])
    is_preprocessing_slide = slide_data['title'] and ('Preparing for Battle' in slide_data['title'] or 'Data Preprocessing' in slide_data['title'])
    is_data_viz_slide = slide_data['title'] and 'Data Visualization' in slide_data['title']
    is_correlations_slide = slide_data['title'] and 'Key Correlations' in slide_data['title']
    is_pca_slide = slide_data['title'] and 'Principal Component Analysis' in slide_data['title']
    is_model_performance_slide = slide_data['title'] and 'Moment of Truth' in slide_data['title']
    is_lda_vs_qda_slide = slide_data['title'] and ('LDA vs QDA' in slide_data['title'] or 'Detailed Results' in slide_data['title'])
    is_blind_spot_slide = slide_data['title'] and 'Critical Blind Spot' in slide_data['title']
    is_plot_twist_slide = slide_data['title'] and 'Plot Twist' in slide_data['title']
    is_dashboard_slide = slide_data['title'] and 'Interactive Dashboard' in slide_data['title']
    is_learned_slide = slide_data['title'] and ('What We Learned' in slide_data['title'] or 'Journey\'s End' in slide_data['title'])
    is_questions_slide = slide_data['title'] and 'Questions' in slide_data['title']
    
    # Special handling for "Mission" slide - two column layout with image on right
    if is_mission_slide and slide_data['images']:
        # Title at top
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.2), 
                            content_width, Inches(0.6),
                            font_size=28, bold=True, 
                            color='#005088', alignment=PP_ALIGN.LEFT)
            y_pos = Inches(1.0)
        
        # Left column: Text content
        left_col_width = Inches(4.5)
        right_col_left = Inches(5.5)
        right_col_width = Inches(4)
        
        # Find the two sections by looking for h4 elements with specific classes
        challenge_header = soup.find('h4', class_='text-primary')
        quest_header = soup.find('h4', class_='text-accent')
        
        challenge_text = ""
        quest_text = ""
        
        # Extract "The Challenge" text
        if challenge_header:
            # Find the paragraph that follows this header
            next_p = challenge_header.find_next_sibling('p')
            if next_p:
                challenge_text = clean_text(next_p.get_text())
            else:
                # Try finding in parent
                parent = challenge_header.parent
                if parent:
                    p_elem = parent.find('p')
                    if p_elem:
                        challenge_text = clean_text(p_elem.get_text())
        
        # Extract "Our Quest" text
        if quest_header:
            # Find the paragraph that follows this header
            next_p = quest_header.find_next_sibling('p')
            if next_p:
                quest_text = clean_text(next_p.get_text())
            else:
                # Try finding in parent
                parent = quest_header.parent
                if parent:
                    p_elem = parent.find('p')
                    if p_elem:
                        quest_text = clean_text(p_elem.get_text())
        
        # Fallback: search all paragraphs
        if not challenge_text or not quest_text:
            paragraphs = soup.find_all('p')
            for para in paragraphs:
                text = clean_text(para.get_text())
                if 'silent killer' in text.lower() or 'irreversible damage' in text.lower():
                    challenge_text = text
                elif 'model' in text.lower() and ('lda' in text.lower() or 'qda' in text.lower()):
                    quest_text = text
        
        # Add "The Challenge" section
        if challenge_text:
            # Section header with icon indicator
            add_text_element(slide, "The Challenge",
                            left_margin, y_pos,
                            left_col_width, Inches(0.4),
                            font_size=18, bold=True,
                            color='#005088', alignment=PP_ALIGN.LEFT)
            y_pos += Inches(0.5)
            
            # Section text
            add_text_element(slide, challenge_text,
                            left_margin, y_pos,
                            left_col_width, Inches(1.0),
                            font_size=11, bold=False,
                            color='#334155', alignment=PP_ALIGN.LEFT)
            y_pos += Inches(1.2)
        
        # Add "Our Quest" section
        if quest_text:
            # Section header
            add_text_element(slide, "Our Quest",
                            left_margin, y_pos,
                            left_col_width, Inches(0.4),
                            font_size=18, bold=True,
                            color='#11caa0', alignment=PP_ALIGN.LEFT)
            y_pos += Inches(0.5)
            
            # Section text
            add_text_element(slide, quest_text,
                            left_margin, y_pos,
                            left_col_width, Inches(1.0),
                            font_size=11, bold=False,
                            color='#334155', alignment=PP_ALIGN.LEFT)
        
        # Right column: Image
        img_data = slide_data['images'][0]
        img_path = img_data['path']
        if os.path.exists(img_path):
            try:
                # Position image on the right side, starting below title
                img_y = Inches(1.2)
                width, height = get_image_size(img_path, right_col_width, Inches(5.5))
                # Make sure image fits in right column
                if width > right_col_width:
                    scale = right_col_width / width
                    width = right_col_width
                    height = height * scale
                
                slide.shapes.add_picture(img_path, right_col_left, img_y, width, height)
            except Exception as e:
                print(f"    Warning: Could not add image {img_path}: {e}")
        
        return
    
    # Special handling for "537 Million" hook slide
    if is_hook_slide:
        # Center everything
        center_left = Inches(1)
        center_width = Inches(8)
        
        # Find the large number in h1
        h1_elem = soup.find('h1')
        if h1_elem:
            number_text = clean_text(h1_elem.get_text())
            if '537' in number_text:
                # Large red centered number
                add_text_element(slide, "537 Million",
                                center_left, Inches(1.8),
                                center_width, Inches(1.5),
                                font_size=80, bold=True,
                                color='#ef4444', alignment=PP_ALIGN.CENTER)
                y_pos = Inches(3.6)
        
        # Add subtitle text - "People worldwide live with diabetes"
        paragraphs = soup.find_all('p')
        for para in paragraphs:
            text = clean_text(para.get_text())
            if not text:
                continue
            
            # Skip source for now, handle it separately
            if 'Source' in text:
                continue
            
            # "People worldwide" - regular text
            if 'People worldwide' in text or 'worldwide' in text.lower():
                add_text_element(slide, text,
                                center_left, y_pos,
                                center_width, Inches(0.4),
                                font_size=14, bold=False,
                                color='#334155', alignment=PP_ALIGN.CENTER)
                y_pos += Inches(0.6)
            
            # "But here's the problem" - bold blue
            elif 'problem' in text.lower() and 'But' in text:
                add_text_element(slide, text,
                                center_left, y_pos,
                                center_width, Inches(0.4),
                                font_size=16, bold=True,
                                color='#005088', alignment=PP_ALIGN.CENTER)
                y_pos += Inches(0.5)
            
            # "Many cases go undetected" - regular gray
            elif 'undetected' in text.lower() or 'complications' in text.lower():
                add_text_element(slide, text,
                                center_left, y_pos,
                                center_width, Inches(0.4),
                                font_size=13, bold=False,
                                color='#64748b', alignment=PP_ALIGN.CENTER)
                y_pos += Inches(0.5)
        
        # Add source in smaller italic text at the end
        for para in paragraphs:
            text = clean_text(para.get_text())
            if 'Source' in text:
                add_text_element(slide, text,
                                center_left, y_pos,
                                center_width, Inches(0.3),
                                font_size=10, bold=False,
                                color='#64748b', alignment=PP_ALIGN.CENTER)
                break
        return
    
    # Special handling for "Our Data" slide - three column cards layout
    if is_data_slide:
        # Title at top
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.2), 
                            content_width, Inches(0.6),
                            font_size=28, bold=True, 
                            color='#005088', alignment=PP_ALIGN.LEFT)
            y_pos = Inches(1.0)
        
        # Subtitle
        subtitle = soup.find('p', string=lambda x: x and 'real-world' in x.lower() if x else False)
        if not subtitle:
            paragraphs = soup.find_all('p')
            for p in paragraphs:
                if 'real-world' in clean_text(p.get_text()).lower():
                    subtitle = p
                    break
        
        if subtitle:
            subtitle_text = clean_text(subtitle.get_text())
            add_text_element(slide, subtitle_text,
                            left_margin, Inches(0.9),
                            content_width, Inches(0.3),
                            font_size=12, bold=False,
                            color='#64748b', alignment=PP_ALIGN.LEFT)
        
        # Three column cards
        card_width = Inches(2.8)
        card_height = Inches(3.5)
        card_gap = Inches(0.3)
        card_y = Inches(1.5)
        
        # Find the three cards from three-col div
        three_col = soup.find(class_='three-col')
        cards = []
        if three_col:
            cards = three_col.find_all('div', class_='card', recursive=False)
            if not cards:
                cards = three_col.find_all('div', recursive=False)
        
        # Extract card content
        card1_text = "768 Women"
        card1_details = ["Pima Indian heritage, aged 21+", "Real patients, real data"]
        card2_text = "8 Predictors"
        card2_details = ["Glucose, BMI, Insulin, Age, Blood Pressure, Pregnancies, Skin Thickness, Pedigree", "Simple, non-invasive measurements"]
        card3_text = "Binary Outcome"
        card3_details = ["0 = No Diabetes", "1 = Diabetes", "Clear classification goal"]
        
        # Try to extract from actual card divs
        if cards and len(cards) >= 3:
            # Card 1
            card1 = cards[0]
            h4 = card1.find('h4')
            if h4:
                card1_text = clean_text(h4.get_text())
            ps = card1.find_all('p')
            card1_details = []
            for p in ps:
                text = clean_text(p.get_text())
                if text and len(text) > 5:
                    card1_details.append(text)
            
            # Card 2
            card2 = cards[1]
            h4 = card2.find('h4')
            if h4:
                card2_text = clean_text(h4.get_text())
            ps = card2.find_all('p')
            card2_details = []
            for p in ps:
                text = clean_text(p.get_text())
                if text and len(text) > 5:
                    card2_details.append(text)
            
            # Card 3
            card3 = cards[2]
            h4 = card3.find('h4')
            if h4:
                card3_text = clean_text(h4.get_text())
            ps = card3.find_all('p')
            card3_details = []
            for p in ps:
                text = clean_text(p.get_text())
                if text and len(text) > 5:
                    # Handle bold text in Binary Outcome
                    if '<b>' in str(p) or '<strong>' in str(p):
                        # Extract bold parts
                        bold_parts = p.find_all(['b', 'strong'])
                        if bold_parts:
                            text = ' = '.join([clean_text(b.get_text()) for b in bold_parts])
                    card3_details.append(text)
        
        # Draw Card 1 (Blue border - left)
        card1_left = left_margin
        if card1_text:
            # Card background (white with blue top border effect)
            add_text_element(slide, card1_text,
                            card1_left, card_y,
                            card_width, Inches(0.5),
                            font_size=20, bold=True,
                            color='#005088', alignment=PP_ALIGN.CENTER)
            detail_y = card_y + Inches(0.6)
            for detail in card1_details[:2]:
                add_text_element(slide, detail,
                                card1_left + Inches(0.1), detail_y,
                                card_width - Inches(0.2), Inches(0.4),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                detail_y += Inches(0.4)
        
        # Draw Card 2 (Green border - middle)
        card2_left = card1_left + card_width + card_gap
        if card2_text:
            add_text_element(slide, card2_text,
                            card2_left, card_y,
                            card_width, Inches(0.5),
                            font_size=20, bold=True,
                            color='#11caa0', alignment=PP_ALIGN.CENTER)
            detail_y = card_y + Inches(0.6)
            for detail in card2_details[:2]:
                add_text_element(slide, detail,
                                card2_left + Inches(0.1), detail_y,
                                card_width - Inches(0.2), Inches(0.4),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                detail_y += Inches(0.4)
        
        # Draw Card 3 (Red border - right)
        card3_left = card2_left + card_width + card_gap
        if card3_text:
            add_text_element(slide, card3_text,
                            card3_left, card_y,
                            card_width, Inches(0.5),
                            font_size=20, bold=True,
                            color='#ef4444', alignment=PP_ALIGN.CENTER)
            detail_y = card_y + Inches(0.6)
            for detail in card3_details[:3]:
                add_text_element(slide, detail,
                                card3_left + Inches(0.1), detail_y,
                                card_width - Inches(0.2), Inches(0.4),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                detail_y += Inches(0.35)
        
        return
    
    # Special handling for "Class Imbalance" slide - chart on left, text boxes on right
    if is_class_imbalance_slide:
        # Title at top
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.2), 
                            content_width, Inches(0.6),
                            font_size=28, bold=True, 
                            color='#005088', alignment=PP_ALIGN.LEFT)
            y_pos = Inches(1.0)
        
        # Subtitle
        subtitle = soup.find('p', string=lambda x: x and 'first challenge' in x.lower() if x else False)
        if not subtitle:
            paragraphs = soup.find_all('p')
            for p in paragraphs:
                if 'first challenge' in clean_text(p.get_text()).lower():
                    subtitle = p
                    break
        
        if subtitle:
            subtitle_text = clean_text(subtitle.get_text())
            add_text_element(slide, subtitle_text,
                            left_margin, Inches(0.9),
                            content_width, Inches(0.3),
                            font_size=11, bold=False,
                            color='#64748b', alignment=PP_ALIGN.LEFT)
        
        # Two column layout
        left_col_width = Inches(5)
        right_col_left = Inches(5.5)
        right_col_width = Inches(4)
        
        # Left column: Image/Chart
        if slide_data['images']:
            img_data = slide_data['images'][0]
            img_path = img_data['path']
            if os.path.exists(img_path):
                try:
                    img_y = Inches(1.3)
                    width, height = get_image_size(img_path, left_col_width, Inches(5))
                    # Make sure image fits
                    if width > left_col_width:
                        scale = left_col_width / width
                        width = left_col_width
                        height = height * scale
                    
                    slide.shapes.add_picture(img_path, left_margin, img_y, width, height)
                except Exception as e:
                    print(f"    Warning: Could not add image {img_path}: {e}")
        
        # Right column: Text boxes
        right_y = Inches(1.3)
        
        # Find "The Numbers" section
        numbers_details = []
        # Look for bg-light div with "The Numbers" header
        bg_light = soup.find(class_='bg-light')
        if bg_light:
            # Find the list items
            for li in bg_light.find_all('li'):
                li_text = clean_text(li.get_text())
                if li_text and ('No Diabetes' in li_text or 'Diabetes' in li_text):
                    numbers_details.append(li_text)
        
        # Fallback: search for list items with the numbers
        if not numbers_details:
            for li in soup.find_all('li'):
                li_text = clean_text(li.get_text())
                if '65%' in li_text or '35%' in li_text:
                    numbers_details.append(li_text)
        
        # Final fallback: use default values
        if not numbers_details:
            numbers_details = ["No Diabetes: ~500 cases (65%)", "Diabetes: ~265 cases (35%)"]
        
        # "The Numbers" box (light blue background)
        if numbers_details:
            # Box header
            add_text_element(slide, "The Numbers",
                            right_col_left, right_y,
                            right_col_width, Inches(0.4),
                            font_size=16, bold=True,
                            color='#005088', alignment=PP_ALIGN.LEFT)
            right_y += Inches(0.5)
            
            # Box content
            for detail in numbers_details[:2]:
                add_text_element(slide, f"• {detail}",
                                right_col_left + Inches(0.1), right_y,
                                right_col_width - Inches(0.2), Inches(0.3),
                                font_size=10, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                right_y += Inches(0.35)
        
        right_y += Inches(0.2)
        
        # Find "Warning Sign" section
        warning_text = ""
        # Look for div with background-color #fef2f2 or border-left
        warning_div = soup.find('div', style=lambda x: x and ('fef2f2' in x.lower() or 'ef4444' in x.lower()) if x else False)
        if warning_div:
            p = warning_div.find('p')
            if p:
                warning_text = clean_text(p.get_text())
                # Remove the "⚠️ Warning Sign:" prefix if present
                if 'Warning Sign' in warning_text:
                    warning_text = warning_text.split(':', 1)[-1].strip()
        
        # Fallback: search for paragraph with warning content
        if not warning_text:
            for p in soup.find_all('p'):
                text = clean_text(p.get_text())
                if '2:1 imbalance' in text.lower() or ('imbalance' in text.lower() and 'majority class' in text.lower()):
                    warning_text = text
                    if 'Warning Sign' in warning_text:
                        warning_text = warning_text.split(':', 1)[-1].strip()
                    break
        
        # Final fallback
        if not warning_text:
            warning_text = "A 2:1 imbalance means our model might favor the majority class. This will become critical later in our story."
        
        # "Warning Sign" box (light orange/red background)
        if warning_text:
            # Box header
            add_text_element(slide, "⚠️ Warning Sign:",
                            right_col_left, right_y,
                            right_col_width, Inches(0.4),
                            font_size=16, bold=True,
                            color='#92400e', alignment=PP_ALIGN.LEFT)
            right_y += Inches(0.5)
            
            # Box content
            add_text_element(slide, warning_text,
                            right_col_left + Inches(0.1), right_y,
                            right_col_width - Inches(0.2), Inches(0.8),
                            font_size=10, bold=False,
                            color='#92400e', alignment=PP_ALIGN.LEFT)
        
        return
    
    # Special handling for "Data Preprocessing" slide - two boxes side by side
    if is_preprocessing_slide:
        # Title at top
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.2), 
                            content_width, Inches(0.6),
                            font_size=28, bold=True, 
                            color='#005088', alignment=PP_ALIGN.LEFT)
            y_pos = Inches(1.0)
        
        # Subtitle
        subtitle = soup.find('p', string=lambda x: x and 'clean our data' in x.lower() if x else False)
        if not subtitle:
            paragraphs = soup.find_all('p')
            for p in paragraphs:
                if 'clean our data' in clean_text(p.get_text()).lower():
                    subtitle = p
                    break
        
        if subtitle:
            subtitle_text = clean_text(subtitle.get_text())
            add_text_element(slide, subtitle_text,
                            left_margin, Inches(0.9),
                            content_width, Inches(0.3),
                            font_size=11, bold=False,
                            color='#64748b', alignment=PP_ALIGN.LEFT)
        
        # Two column layout for the two steps
        box_width = Inches(4.3)
        box_height = Inches(2.5)
        box_gap = Inches(0.4)
        box_y = Inches(1.4)
        
        # Find the two bg-light divs (steps)
        bg_lights = soup.find_all(class_='bg-light')
        step1_text = ""
        step1_motto = ""
        step2_text = ""
        step2_motto = ""
        
        if len(bg_lights) >= 2:
            # Step 1: Data Cleaning
            step1_div = bg_lights[0]
            h4 = step1_div.find('h4')
            if h4:
                step1_title = clean_text(h4.get_text())
            ps = step1_div.find_all('p')
            for p in ps:
                text = clean_text(p.get_text())
                # Check if it's the main description (longer text)
                if text and len(text) > 30 and ('impossible zero' in text.lower() or 'median values' in text.lower() or 'NA' in text):
                    step1_text = text
                # Check if it's the motto (shorter, italicized)
                elif text and ('Quality data' in text or 'reliable results' in text.lower()):
                    step1_motto = text
        
            # Step 2: Train-Test Split
            step2_div = bg_lights[1]
            h4 = step2_div.find('h4')
            if h4:
                step2_title = clean_text(h4.get_text())
            ps = step2_div.find_all('p')
            for p in ps:
                text = clean_text(p.get_text())
                # Check if it's the main description (longer text)
                if text and len(text) > 20 and ('70%' in text or 'Training' in text or 'Testing' in text or 'stratified' in text.lower()):
                    step2_text = text
                # Check if it's the motto (shorter, italicized)
                elif text and ('Fair evaluation' in text or 'fair splits' in text.lower()):
                    step2_motto = text
        
        # Fallback: extract from all paragraphs
        if not step1_text or not step2_text:
            paragraphs = soup.find_all('p')
            for p in paragraphs:
                text = clean_text(p.get_text())
                if 'impossible zero' in text.lower() or 'median values' in text.lower():
                    step1_text = text
                elif '70%' in text or ('Training' in text and 'Testing' in text):
                    step2_text = text
                elif 'Quality data' in text:
                    step1_motto = text
                elif 'Fair evaluation' in text:
                    step2_motto = text
        
        # Draw Step 1 box (left, blue border)
        step1_left = left_margin
        if step1_text:
            # Box title
            add_text_element(slide, "Step 1: Data Cleaning",
                            step1_left + Inches(0.2), box_y,
                            box_width - Inches(0.4), Inches(0.4),
                            font_size=16, bold=True,
                            color='#005088', alignment=PP_ALIGN.LEFT)
            
            # Box description
            add_text_element(slide, step1_text,
                            step1_left + Inches(0.2), box_y + Inches(0.5),
                            box_width - Inches(0.4), Inches(1.2),
                            font_size=10, bold=False,
                            color='#334155', alignment=PP_ALIGN.LEFT)
            
            # Box motto
            if step1_motto:
                add_text_element(slide, step1_motto,
                                step1_left + Inches(0.2), box_y + Inches(1.8),
                                box_width - Inches(0.4), Inches(0.3),
                                font_size=9, bold=False,
                                color='#64748b', alignment=PP_ALIGN.LEFT)
        
        # Draw Step 2 box (right, teal border)
        step2_left = step1_left + box_width + box_gap
        if step2_text:
            # Box title
            add_text_element(slide, "Step 2: Train-Test Split",
                            step2_left + Inches(0.2), box_y,
                            box_width - Inches(0.4), Inches(0.4),
                            font_size=16, bold=True,
                            color='#11caa0', alignment=PP_ALIGN.LEFT)
            
            # Box description
            add_text_element(slide, step2_text,
                            step2_left + Inches(0.2), box_y + Inches(0.5),
                            box_width - Inches(0.4), Inches(1.2),
                            font_size=10, bold=False,
                            color='#334155', alignment=PP_ALIGN.LEFT)
            
            # Box motto
            if step2_motto:
                add_text_element(slide, step2_motto,
                                step2_left + Inches(0.2), box_y + Inches(1.8),
                                box_width - Inches(0.4), Inches(0.3),
                                font_size=9, bold=False,
                                color='#64748b', alignment=PP_ALIGN.LEFT)
        
        return
    
    # Special handling for "Data Visualization" slide - larger image with text below
    if is_data_viz_slide:
        # Title at top
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.15), 
                            content_width, Inches(0.5),
                            font_size=26, bold=True, 
                            color='#005088', alignment=PP_ALIGN.LEFT)
        
        # Image centered - smaller size to leave room for text below
        img_y = Inches(0.7)
        img_height = Inches(3.5)  # Reduced to leave room for caption and insights
        caption_y = None
        
        if slide_data['images']:
            img_data = slide_data['images'][0]
            img_path = img_data['path']
            if os.path.exists(img_path):
                try:
                    # Center image horizontally, smaller to leave room for insights
                    width, height = get_image_size(img_path, Inches(7.5), Inches(3.5))
                    img_left = (Inches(10) - width) / 2
                    slide.shapes.add_picture(img_path, img_left, img_y, width, height)
                    img_height = height
                    
                    # Caption below image with proper spacing
                    caption_y = img_y + height + Inches(0.12)
                    caption = soup.find(class_='image-caption')
                    if caption:
                        caption_text = clean_text(caption.get_text())
                        add_text_element(slide, caption_text,
                                        left_margin, caption_y,
                                        content_width, Inches(0.25),
                                        font_size=9, bold=False,
                                        color='#64748b', alignment=PP_ALIGN.CENTER)
                        caption_y = caption_y + Inches(0.4)  # Position after caption with spacing
                except Exception as e:
                    print(f"    Warning: Could not add image {img_path}: {e}")
                    import traceback
                    traceback.print_exc()
        
        # Key insights box at bottom - ALWAYS add this with proper spacing
        # Calculate position after image and caption
        if caption_y and hasattr(caption_y, 'inches'):
            insights_y = caption_y + Inches(0.1)  # Space after caption
        else:
            # If no caption, position after image with spacing
            insights_y = img_y + img_height + Inches(0.35)
        
        # Ensure insights are visible and don't overlap - use fixed position
        if hasattr(insights_y, 'inches'):
            if insights_y.inches > 5.2:
                insights_y = Inches(4.8)  # Fixed position to prevent overlap
        else:
            # If it's not an Inches object, use fixed position
            insights_y = Inches(4.8)
        
        # Default insights (always use these)
        insights_list = [
            "Glucose appears to be the strongest distinguishing factor (diabetes cases cluster at Glucose > 120-130)",
            "Higher BMI and Age show increased diabetes prevalence, but with significant overlap",
            "No clear linear separation, suggesting non-linear boundaries may be needed"
        ]
        
        # Try to extract from HTML if possible
        insights_div = None
        # Method 1: Look for div with background-color #f8fafc
        for div in soup.find_all('div'):
            style = div.get('style', '')
            if 'f8fafc' in style.lower() or ('background-color' in style.lower() and 'f8' in style.lower()):
                insights_div = div
                break
        
        # Method 2: Look for div containing "Key Insights"
        if not insights_div:
            for div in soup.find_all('div'):
                if 'Key Insights' in div.get_text():
                    insights_div = div
                    break
        
        # Extract from HTML if found
        if insights_div:
            extracted_list = []
            for li in insights_div.find_all('li'):
                li_text = clean_text(li.get_text())
                if li_text and len(li_text) > 10:
                    extracted_list.append(li_text)
            if extracted_list:
                insights_list = extracted_list[:3]
        
        # Always add Key Insights section
        print(f"    Adding Key Insights at y={insights_y.inches if hasattr(insights_y, 'inches') else insights_y}")
        print(f"    Adding {len(insights_list)} insights")
        
        # Add "Key Insights:" header (bold, blue)
        add_text_element(slide, "Key Insights:",
                        left_margin, insights_y,
                        content_width, Inches(0.25),
                        font_size=10, bold=True,
                        color='#005088', alignment=PP_ALIGN.LEFT)
        insights_y += Inches(0.3)
        
        # Add insights as bullet points (matching HTML font sizes)
        for insight in insights_list[:3]:  # Limit to 3 insights
            if hasattr(insights_y, 'inches'):
                if insights_y.inches < 6.8:
                    add_text_element(slide, f"• {insight}",
                                    left_margin + Inches(0.2), insights_y,
                                    content_width - Inches(0.2), Inches(0.4),
                                    font_size=9, bold=False,
                                    color='#334155', alignment=PP_ALIGN.LEFT)
                    insights_y += Inches(0.4)
                    print(f"      Added insight: {insight[:50]}...")
            else:
                # Fallback if insights_y is not an Inches object
                add_text_element(slide, f"• {insight}",
                                left_margin + Inches(0.2), insights_y,
                                content_width - Inches(0.2), Inches(0.4),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                insights_y += Inches(0.4)
                print(f"      Added insight: {insight[:50]}...")
        
        return
    
    # Special handling for "Key Correlations" slide - two column layout
    if is_correlations_slide:
        # Title at top
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.2), 
                            content_width, Inches(0.6),
                            font_size=28, bold=True, 
                            color='#005088', alignment=PP_ALIGN.LEFT)
        
        # Two column layout
        left_col_width = Inches(4.5)
        right_col_left = Inches(5.5)
        right_col_width = Inches(4)
        col_y = Inches(1.0)
        
        # Left column: Strong and Weak Correlations
        # Find "Strong Correlations" section
        strong_correlations = []
        weak_correlations = []
        
        # Look for h4 elements with "Strong" or "Weak"
        for h4 in soup.find_all('h4'):
            text = clean_text(h4.get_text())
            parent = h4.parent if h4.parent else h4
            if 'Strong' in text:
                # Find list items under this heading
                for li in parent.find_all('li', recursive=True):
                    li_text = clean_text(li.get_text())
                    if li_text and len(li_text) > 10:
                        strong_correlations.append(li_text)
            elif 'Weak' in text:
                # Find list items under this heading
                for li in parent.find_all('li', recursive=True):
                    li_text = clean_text(li.get_text())
                    if li_text and len(li_text) > 10:
                        weak_correlations.append(li_text)
        
        # Fallback: extract from all list items
        if not strong_correlations or not weak_correlations:
            all_lis = soup.find_all('li')
            for li in all_lis:
                text = clean_text(li.get_text())
                if '0.54' in text or '0.42' in text:
                    strong_correlations.append(text)
                elif '0.00' in text or '0.02' in text:
                    weak_correlations.append(text)
        
        # Default values if not found
        if not strong_correlations:
            strong_correlations = [
                "Age & Pregnancies (0.54): Strongest correlation; older women tend to have higher pregnancy counts.",
                "BMI & Skin Thickness (0.54): Expected physiological relationship - both measure body composition.",
                "Glucose & Insulin (0.42): Moderate positive correlation - insulin response to glucose levels."
            ]
        if not weak_correlations:
            weak_correlations = [
                "BloodPressure & DiabetesPedigreeFunction (0.00): No linear relationship.",
                "Pregnancies & BMI (0.02): Nearly independent."
            ]
        
        # Add "Strong Correlations" section
        add_text_element(slide, "Strong Correlations",
                        left_margin, col_y,
                        left_col_width, Inches(0.35),
                        font_size=16, bold=True,
                        color='#005088', alignment=PP_ALIGN.LEFT)
        col_y += Inches(0.4)
        
        for corr in strong_correlations[:3]:
            if col_y < Inches(4.5):
                add_text_element(slide, f"• {corr}",
                                left_margin + Inches(0.2), col_y,
                                left_col_width - Inches(0.2), Inches(0.35),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                col_y += Inches(0.4)
        
        col_y += Inches(0.2)
        
        # Add "Weak Correlations" section
        add_text_element(slide, "Weak Correlations",
                        left_margin, col_y,
                        left_col_width, Inches(0.35),
                        font_size=16, bold=True,
                        color='#005088', alignment=PP_ALIGN.LEFT)
        col_y += Inches(0.4)
        
        for corr in weak_correlations[:2]:
            if col_y < Inches(5.5):
                add_text_element(slide, f"• {corr}",
                                left_margin + Inches(0.2), col_y,
                                left_col_width - Inches(0.2), Inches(0.35),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                col_y += Inches(0.4)
        
        # Right column: Implications box
        implications_y = Inches(1.0)
        
        # Find implications text
        implications_text = []
        bg_light = soup.find(class_='bg-light')
        if bg_light:
            for p in bg_light.find_all('p'):
                text = clean_text(p.get_text())
                if text and len(text) > 20:
                    implications_text.append(text)
        
        # Default implications
        if not implications_text:
            implications_text = [
                "No perfect multicollinearity: All correlations are moderate, suggesting each predictor contributes unique information. This supports the use of discriminant analysis methods.",
                "No single dominant predictor: The outcome cannot be predicted by a single variable, necessitating multivariate approaches like LDA/QDA."
            ]
        
        # Add "Implications" header
        add_text_element(slide, "Implications",
                        right_col_left, implications_y,
                        right_col_width, Inches(0.35),
                        font_size=16, bold=True,
                        color='#11caa0', alignment=PP_ALIGN.LEFT)
        implications_y += Inches(0.4)
        
        # Add implications text
        for impl in implications_text[:2]:
            if implications_y < Inches(6.0):
                add_text_element(slide, impl,
                                right_col_left + Inches(0.1), implications_y,
                                right_col_width - Inches(0.2), Inches(0.6),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                implications_y += Inches(0.7)
        
        return
    
    # Special handling for "Principal Component Analysis" slide - image left, cards right
    if is_pca_slide:
        # Title at top
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.2), 
                            content_width, Inches(0.6),
                            font_size=28, bold=True, 
                            color='#005088', alignment=PP_ALIGN.LEFT)
        
        # Two column layout
        left_col_width = Inches(5)
        right_col_left = Inches(5.5)
        right_col_width = Inches(4)
        
        # Left column: Image
        img_y = Inches(1.0)
        if slide_data['images']:
            img_data = slide_data['images'][0]
            img_path = img_data['path']
            if os.path.exists(img_path):
                try:
                    width, height = get_image_size(img_path, left_col_width, Inches(4.5))
                    slide.shapes.add_picture(img_path, left_margin, img_y, width, height)
                    
                    # Caption below image
                    caption = soup.find(class_='image-caption')
                    if caption:
                        caption_text = clean_text(caption.get_text())
                        add_text_element(slide, caption_text,
                                        left_margin, img_y + height + Inches(0.05),
                                        left_col_width, Inches(0.25),
                                        font_size=8, bold=False,
                                        color='#64748b', alignment=PP_ALIGN.LEFT)
                except Exception as e:
                    print(f"    Warning: Could not add image {img_path}: {e}")
        
        # Right column: Three cards
        card_y = Inches(1.0)
        card_height = Inches(1.0)
        card_gap = Inches(0.15)
        
        # Find the three cards
        cards = soup.find_all(class_='card')
        if not cards:
            # Look for cards in two-col div
            two_col = soup.find(class_='two-col')
            if two_col:
                cards = two_col.find_all('div', class_='card', recursive=False)
        
        # Card 1: 47.23% Variance Explained
        card1_text = ""
        card1_details = ""
        if cards and len(cards) >= 1:
            card1 = cards[0]
            h4 = card1.find('h4')
            if h4:
                card1_text = clean_text(h4.get_text())
            p = card1.find('p')
            if p:
                card1_details = clean_text(p.get_text())
        
        # Fallback
        if not card1_text:
            card1_text = "47.23% Variance Explained"
            card1_details = "First 2 components (PC1 + PC2) capture less than half of total variance."
        
        # Card 2: 5 Components Needed
        card2_text = ""
        card2_details = ""
        if cards and len(cards) >= 2:
            card2 = cards[1]
            h4 = card2.find('h4')
            if h4:
                card2_text = clean_text(h4.get_text())
            p = card2.find('p')
            if p:
                card2_details = clean_text(p.get_text())
        
        # Fallback
        if not card2_text:
            card2_text = "5 Components Needed"
            card2_details = "Required to capture >80% of the dataset's information."
        
        # Card 3: Key Observation
        card3_text = ""
        card3_details = ""
        bg_light = soup.find(class_='bg-light')
        if bg_light:
            h4 = bg_light.find('h4')
            if h4:
                card3_text = clean_text(h4.get_text())
            p = bg_light.find('p')
            if p:
                card3_details = clean_text(p.get_text())
        
        # Fallback
        if not card3_text:
            card3_text = "Key Observation"
            card3_details = "Heavy Overlap: The first two principal components show significant overlap between diabetes and no-diabetes groups, indicating that linear separation in reduced space is challenging. This supports the need for discriminant analysis methods."
        
        # Draw Card 1 (teal/accent color)
        if card1_text:
            add_text_element(slide, card1_text,
                            right_col_left + Inches(0.1), card_y,
                            right_col_width - Inches(0.2), Inches(0.35),
                            font_size=14, bold=True,
                            color='#11caa0', alignment=PP_ALIGN.LEFT)
            if card1_details:
                add_text_element(slide, card1_details,
                                right_col_left + Inches(0.1), card_y + Inches(0.4),
                                right_col_width - Inches(0.2), Inches(0.6),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
        
        card_y += card_height + card_gap
        
        # Draw Card 2
        if card2_text:
            add_text_element(slide, card2_text,
                            right_col_left + Inches(0.1), card_y,
                            right_col_width - Inches(0.2), Inches(0.35),
                            font_size=14, bold=True,
                            color='#005088', alignment=PP_ALIGN.LEFT)
            if card2_details:
                add_text_element(slide, card2_details,
                                right_col_left + Inches(0.1), card_y + Inches(0.4),
                                right_col_width - Inches(0.2), Inches(0.6),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
        
        card_y += card_height + card_gap
        
        # Draw Card 3 (Key Observation)
        if card3_text:
            add_text_element(slide, card3_text,
                            right_col_left + Inches(0.1), card_y,
                            right_col_width - Inches(0.2), Inches(0.35),
                            font_size=14, bold=True,
                            color='#005088', alignment=PP_ALIGN.LEFT)
            if card3_details:
                # Truncate if too long
                if len(card3_details) > 150:
                    card3_details = card3_details[:147] + "..."
                add_text_element(slide, card3_details,
                                right_col_left + Inches(0.1), card_y + Inches(0.4),
                                right_col_width - Inches(0.2), Inches(0.8),
                                font_size=8, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
        
        return
    
    # Special handling for "The Moment of Truth: Model Performance" slide - table left, ROC curve right
    if is_model_performance_slide:
        # Title at top
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.2), 
                            content_width, Inches(0.6),
                            font_size=28, bold=True, 
                            color='#005088', alignment=PP_ALIGN.LEFT)
        
        # Subtitle
        subtitle = soup.find('p', class_='subtitle')
        if not subtitle:
            # Look for any p tag with subtitle-like text
            for p in soup.find_all('p'):
                text = clean_text(p.get_text())
                if 'rigorous testing' in text.lower() or 'results are in' in text.lower():
                    subtitle = p
                    break
        
        if subtitle:
            subtitle_text = clean_text(subtitle.get_text())
            add_text_element(slide, subtitle_text,
                            left_margin, Inches(0.8),
                            content_width, Inches(0.3),
                            font_size=12, bold=False,
                            color='#64748b', alignment=PP_ALIGN.LEFT)
        
        # Two column layout
        left_col_width = Inches(4.5)
        right_col_left = Inches(5.5)
        right_col_width = Inches(4)
        col_y = Inches(1.3)
        
        # Left column: Performance metrics table
        # Find table
        table = soup.find('table')
        if table:
            # Table header
            headers = []
            thead = table.find('thead')
            if thead:
                for th in thead.find_all('th'):
                    headers.append(clean_text(th.get_text()))
            else:
                # Try first row
                first_row = table.find('tr')
                if first_row:
                    for th in first_row.find_all(['th', 'td']):
                        headers.append(clean_text(th.get_text()))
            
            # Default headers if not found
            if not headers:
                headers = ['Metric', 'LDA (Linear)', 'QDA (Quadratic)']
            
            # Table rows
            rows = []
            tbody = table.find('tbody')
            if tbody:
                for tr in tbody.find_all('tr'):
                    row_data = []
                    for td in tr.find_all('td'):
                        cell_text = clean_text(td.get_text())
                        # Check for green/highlighted text
                        if td.find(class_='text-success') or 'text-success' in str(td.get('class', [])):
                            cell_text = f"**{cell_text}**"  # Mark as bold for green
                        row_data.append(cell_text)
                    if row_data:
                        rows.append(row_data)
            else:
                # Try all rows except first
                all_rows = table.find_all('tr')
                for tr in all_rows[1:]:
                    row_data = []
                    for td in tr.find_all('td'):
                        cell_text = clean_text(td.get_text())
                        if td.find(class_='text-success') or 'text-success' in str(td.get('class', [])):
                            cell_text = f"**{cell_text}**"
                        row_data.append(cell_text)
                    if row_data:
                        rows.append(row_data)
            
            # Default rows if not found
            if not rows:
                rows = [
                    ['Accuracy', '76.09%', '71.30%'],
                    ['Sensitivity (Healthy)', '86.67%', '82.00%'],
                    ['Specificity (Diabetic)', '56.25%', '51.25%'],
                    ['AUC (ROC)', '0.835', '0.806']
                ]
            
            # Draw table
            table_y = col_y
            row_height = Inches(0.4)
            col_widths = [Inches(2.2), Inches(1.1), Inches(1.1)]
            
            # Header row
            header_x = left_margin
            for i, header in enumerate(headers[:3]):
                add_text_element(slide, header,
                                header_x, table_y,
                                col_widths[i], row_height,
                                font_size=10, bold=True,
                                color='#005088', alignment=PP_ALIGN.LEFT)
                header_x += col_widths[i]
            table_y += row_height
            
            # Data rows
            for row in rows[:4]:
                row_x = left_margin
                for i, cell in enumerate(row[:3]):
                    # Check if cell should be green/bold
                    is_highlighted = cell.startswith('**') and cell.endswith('**')
                    if is_highlighted:
                        cell = cell.strip('*')
                    
                    add_text_element(slide, cell,
                                    row_x, table_y,
                                    col_widths[i], row_height,
                                    font_size=9, bold=is_highlighted,
                                    color='#10b981' if is_highlighted else '#334155',
                                    alignment=PP_ALIGN.LEFT)
                    row_x += col_widths[i]
                table_y += row_height
        
        # Right column: ROC curve image
        if slide_data['images']:
            img_data = slide_data['images'][0]
            img_path = img_data['path']
            if os.path.exists(img_path):
                try:
                    width, height = get_image_size(img_path, right_col_width, Inches(4.5))
                    # Center vertically with table
                    img_y = col_y
                    slide.shapes.add_picture(img_path, right_col_left, img_y, width, height)
                except Exception as e:
                    print(f"    Warning: Could not add ROC image {img_path}: {e}")
        
        return
    
    # Special handling for "LDA vs QDA: Detailed Results" slide - two cards side by side
    if is_lda_vs_qda_slide:
        # Title at top
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.2), 
                            content_width, Inches(0.6),
                            font_size=28, bold=True, 
                            color='#005088', alignment=PP_ALIGN.LEFT)
        
        # Two column layout for cards
        left_card_left = left_margin
        left_card_width = Inches(4.5)
        right_card_left = Inches(5.5)
        right_card_width = Inches(4)
        card_y = Inches(1.0)
        card_height = Inches(4.5)
        
        # Find the two cards
        cards = soup.find_all(class_='card')
        if not cards:
            # Look in two-col div
            two_col = soup.find(class_='two-col')
            if two_col:
                cards = two_col.find_all('div', class_='card', recursive=False)
        
        # Left card: LDA
        lda_card = None
        if cards and len(cards) >= 1:
            lda_card = cards[0]
        
        # Extract LDA metrics
        lda_metrics = {}
        lda_strength = ""
        if lda_card:
            # Find all list items or paragraphs with metrics
            for li in lda_card.find_all('li'):
                text = clean_text(li.get_text())
                if 'Accuracy' in text:
                    lda_metrics['Accuracy'] = text.split(':')[-1].strip()
                elif 'Sensitivity' in text:
                    lda_metrics['Sensitivity'] = text.split(':')[-1].strip()
                elif 'Specificity' in text:
                    lda_metrics['Specificity'] = text.split(':')[-1].strip()
                elif 'AUC-ROC' in text or 'AUC' in text:
                    lda_metrics['AUC-ROC'] = text.split(':')[-1].strip()
                elif 'Precision' in text:
                    lda_metrics['Precision'] = text.split(':')[-1].strip()
                elif 'F1-Score' in text or 'F1' in text:
                    lda_metrics['F1-Score'] = text.split(':')[-1].strip()
            
            # Find Key Strength
            bg_light = lda_card.find(class_='bg-light')
            if bg_light:
                p = bg_light.find('p')
                if p:
                    lda_strength = clean_text(p.get_text())
        
        # Default LDA values
        if not lda_metrics:
            lda_metrics = {
                'Accuracy': '76.09%',
                'Sensitivity': '86.67%',
                'Specificity': '56.25%',
                'AUC-ROC': '0.835',
                'Precision': '0.71',
                'F1-Score': '0.78'
            }
        if not lda_strength:
            lda_strength = "More robust to class imbalance, better generalization."
        
        # Right card: QDA
        qda_card = None
        if cards and len(cards) >= 2:
            qda_card = cards[1]
        
        # Extract QDA metrics
        qda_metrics = {}
        qda_limitation = ""
        if qda_card:
            for li in qda_card.find_all('li'):
                text = clean_text(li.get_text())
                if 'Accuracy' in text:
                    qda_metrics['Accuracy'] = text.split(':')[-1].strip()
                elif 'Sensitivity' in text:
                    qda_metrics['Sensitivity'] = text.split(':')[-1].strip()
                elif 'Specificity' in text:
                    qda_metrics['Specificity'] = text.split(':')[-1].strip()
                elif 'AUC-ROC' in text or 'AUC' in text:
                    qda_metrics['AUC-ROC'] = text.split(':')[-1].strip()
                elif 'Precision' in text:
                    qda_metrics['Precision'] = text.split(':')[-1].strip()
                elif 'F1-Score' in text or 'F1' in text:
                    qda_metrics['F1-Score'] = text.split(':')[-1].strip()
            
            # Find Limitation
            bg_light = qda_card.find(class_='bg-light')
            if bg_light:
                p = bg_light.find('p')
                if p:
                    qda_limitation = clean_text(p.get_text())
        
        # Default QDA values
        if not qda_metrics:
            qda_metrics = {
                'Accuracy': '71.30%',
                'Sensitivity': '82.00%',
                'Specificity': '51.25%',
                'AUC-ROC': '0.806',
                'Precision': '0.66',
                'F1-Score': '0.73'
            }
        if not qda_limitation:
            qda_limitation = "Overfits on smaller class, less stable with imbalanced data."
        
        # Draw LDA card
        lda_y = card_y
        add_text_element(slide, "LDA (Linear Discriminant Analysis)",
                        left_card_left + Inches(0.1), lda_y,
                        left_card_width - Inches(0.2), Inches(0.35),
                        font_size=14, bold=True,
                        color='#005088', alignment=PP_ALIGN.LEFT)
        lda_y += Inches(0.45)
        
        # LDA metrics (all in green)
        metric_order = ['Accuracy', 'Sensitivity', 'Specificity', 'AUC-ROC', 'Precision', 'F1-Score']
        for metric in metric_order:
            if metric in lda_metrics:
                value = lda_metrics[metric]
                add_text_element(slide, f"{metric}: {value}",
                                left_card_left + Inches(0.15), lda_y,
                                left_card_width - Inches(0.3), Inches(0.3),
                                font_size=9, bold=False,
                                color='#10b981', alignment=PP_ALIGN.LEFT)
                lda_y += Inches(0.35)
        
        lda_y += Inches(0.1)
        
        # LDA Key Strength box
        add_text_element(slide, "Key Strength:",
                        left_card_left + Inches(0.1), lda_y,
                        left_card_width - Inches(0.2), Inches(0.25),
                        font_size=9, bold=True,
                        color='#334155', alignment=PP_ALIGN.LEFT)
        lda_y += Inches(0.3)
        add_text_element(slide, lda_strength,
                        left_card_left + Inches(0.1), lda_y,
                        left_card_width - Inches(0.2), Inches(0.5),
                        font_size=8, bold=False,
                        color='#334155', alignment=PP_ALIGN.LEFT)
        
        # Draw QDA card
        qda_y = card_y
        add_text_element(slide, "QDA (Quadratic Discriminant Analysis)",
                        right_card_left + Inches(0.1), qda_y,
                        right_card_width - Inches(0.2), Inches(0.35),
                        font_size=14, bold=True,
                        color='#11caa0', alignment=PP_ALIGN.LEFT)
        qda_y += Inches(0.45)
        
        # QDA metrics (regular color)
        for metric in metric_order:
            if metric in qda_metrics:
                value = qda_metrics[metric]
                add_text_element(slide, f"{metric}: {value}",
                                right_card_left + Inches(0.15), qda_y,
                                right_card_width - Inches(0.3), Inches(0.3),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                qda_y += Inches(0.35)
        
        qda_y += Inches(0.1)
        
        # QDA Limitation box
        add_text_element(slide, "Limitation:",
                        right_card_left + Inches(0.1), qda_y,
                        right_card_width - Inches(0.2), Inches(0.25),
                        font_size=9, bold=True,
                        color='#334155', alignment=PP_ALIGN.LEFT)
        qda_y += Inches(0.3)
        add_text_element(slide, qda_limitation,
                        right_card_left + Inches(0.1), qda_y,
                        right_card_width - Inches(0.2), Inches(0.5),
                        font_size=8, bold=False,
                        color='#334155', alignment=PP_ALIGN.LEFT)
        
        # Winner box at bottom
        winner_y = Inches(5.8)
        winner_text = ""
        
        # Find winner box
        bg_warning = soup.find(class_='bg-warning')
        if bg_warning:
            winner_text = clean_text(bg_warning.get_text())
        else:
            # Look for any div with "Winner" text
            for div in soup.find_all('div'):
                text = clean_text(div.get_text())
                if 'Winner' in text and 'LDA' in text:
                    winner_text = text
                    break
        
        if not winner_text:
            winner_text = "Winner: LDA - Despite theoretical preference for QDA (unequal covariance), LDA demonstrates superior performance with 4.79% higher accuracy and better stability on imbalanced data."
        
        add_text_element(slide, winner_text,
                        left_margin + Inches(0.1), winner_y,
                        content_width - Inches(0.2), Inches(0.6),
                        font_size=9, bold=True,
                        color='#92400e', alignment=PP_ALIGN.LEFT)
        
        return
    
    # Special handling for "The Critical Blind Spot" slide - two cards and critical issue box
    if is_blind_spot_slide:
        # Set background color to light pink
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(254, 242, 242)  # #fef2f2
        
        # Title at top (red)
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.2), 
                            content_width, Inches(0.6),
                            font_size=28, bold=True, 
                            color='#ef4444', alignment=PP_ALIGN.LEFT)
        
        # Subtitle
        subtitle = soup.find('p')
        if subtitle:
            subtitle_text = clean_text(subtitle.get_text())
            if 'sobering' in subtitle_text.lower() or 'reality' in subtitle_text.lower():
                add_text_element(slide, subtitle_text,
                                left_margin, Inches(0.8),
                                content_width, Inches(0.3),
                                font_size=12, bold=False,
                                color='#64748b', alignment=PP_ALIGN.LEFT)
        
        # Two cards side by side
        card_y = Inches(1.4)
        card_width = Inches(4.2)
        card_height = Inches(2.0)
        left_card_left = left_margin
        right_card_left = Inches(5.8)
        
        # Find the two cards
        cards = soup.find_all(class_='card')
        if not cards:
            # Look for divs with card-like structure
            for div in soup.find_all('div'):
                if 'card' in str(div.get('class', [])) or 'width: 250px' in str(div.get('style', '')):
                    cards.append(div)
        
        # Left card: 56% Specificity
        card1_number = "56%"
        card1_label = "Specificity"
        card1_text = "Only 56% of diabetic cases correctly identified"
        
        if cards and len(cards) >= 1:
            card1 = cards[0]
            h1 = card1.find('h1')
            if h1:
                card1_number = clean_text(h1.get_text())
            p_tags = card1.find_all('p')
            if len(p_tags) >= 1:
                card1_label = clean_text(p_tags[0].get_text())
            if len(p_tags) >= 2:
                card1_text = clean_text(p_tags[1].get_text())
        
        # Right card: 44% Missed Cases
        card2_number = "44%"
        card2_label = "Missed Cases"
        card2_text = "44% of diabetic patients go undetected"
        
        if cards and len(cards) >= 2:
            card2 = cards[1]
            h1 = card2.find('h1')
            if h1:
                card2_number = clean_text(h1.get_text())
            p_tags = card2.find_all('p')
            if len(p_tags) >= 1:
                card2_label = clean_text(p_tags[0].get_text())
            if len(p_tags) >= 2:
                card2_text = clean_text(p_tags[1].get_text())
        
        # Draw left card
        card1_y = card_y
        add_text_element(slide, card1_number,
                        left_card_left + Inches(0.2), card1_y,
                        card_width - Inches(0.4), Inches(0.8),
                        font_size=48, bold=True,
                        color='#ef4444', alignment=PP_ALIGN.CENTER)
        card1_y += Inches(0.85)
        add_text_element(slide, card1_label,
                        left_card_left + Inches(0.2), card1_y,
                        card_width - Inches(0.4), Inches(0.3),
                        font_size=11, bold=False,
                        color='#64748b', alignment=PP_ALIGN.CENTER)
        card1_y += Inches(0.35)
        add_text_element(slide, card1_text,
                        left_card_left + Inches(0.2), card1_y,
                        card_width - Inches(0.4), Inches(0.5),
                        font_size=9, bold=False,
                        color='#334155', alignment=PP_ALIGN.CENTER)
        
        # Draw right card
        card2_y = card_y
        add_text_element(slide, card2_number,
                        right_card_left + Inches(0.2), card2_y,
                        card_width - Inches(0.4), Inches(0.8),
                        font_size=48, bold=True,
                        color='#ef4444', alignment=PP_ALIGN.CENTER)
        card2_y += Inches(0.85)
        add_text_element(slide, card2_label,
                        right_card_left + Inches(0.2), card2_y,
                        card_width - Inches(0.4), Inches(0.3),
                        font_size=11, bold=False,
                        color='#64748b', alignment=PP_ALIGN.CENTER)
        card2_y += Inches(0.35)
        add_text_element(slide, card2_text,
                        right_card_left + Inches(0.2), card2_y,
                        card_width - Inches(0.4), Inches(0.5),
                        font_size=9, bold=False,
                        color='#334155', alignment=PP_ALIGN.CENTER)
        
        # Critical Issue box at bottom
        critical_y = Inches(4.0)
        critical_text = ""
        
        # Find critical issue box
        bg_light = soup.find('div', style=lambda x: x and 'background-color' in x and 'fef2f2' in x.lower())
        if not bg_light:
            # Look for div with "Critical Issue" text
            for div in soup.find_all('div'):
                text = clean_text(div.get_text())
                if 'Critical Issue' in text or 'critical issue' in text.lower():
                    critical_text = text
                    break
        
        if not critical_text:
            critical_text = "Critical Issue: The model is excellent at identifying healthy individuals (86.67% sensitivity), but struggles with positive cases due to class imbalance (65% No Diabetes vs 35% Diabetes). This is a significant limitation for clinical application where missing diabetic cases has serious consequences."
        
        # Split into "Critical Issue:" header and body
        if 'Critical Issue:' in critical_text:
            parts = critical_text.split(':', 1)
            header = parts[0] + ":"
            body = parts[1].strip() if len(parts) > 1 else ""
        else:
            header = "Critical Issue:"
            body = critical_text
        
        add_text_element(slide, header,
                        left_margin + Inches(0.1), critical_y,
                        content_width - Inches(0.2), Inches(0.3),
                        font_size=10, bold=True,
                        color='#334155', alignment=PP_ALIGN.LEFT)
        critical_y += Inches(0.35)
        add_text_element(slide, body,
                        left_margin + Inches(0.1), critical_y,
                        content_width - Inches(0.2), Inches(1.2),
                        font_size=9, bold=False,
                        color='#334155', alignment=PP_ALIGN.LEFT)
        
        return
    
    # Special handling for "The Plot Twist: Theory vs. Reality" slide - two column cards
    if is_plot_twist_slide:
        # Title at top
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.2), 
                            content_width, Inches(0.6),
                            font_size=28, bold=True, 
                            color='#005088', alignment=PP_ALIGN.LEFT)
        
        # Subtitle
        subtitle = soup.find('p')
        if subtitle:
            subtitle_text = clean_text(subtitle.get_text())
            if 'data tells' in subtitle_text.lower() or 'textbooks' in subtitle_text.lower():
                add_text_element(slide, subtitle_text,
                                left_margin, Inches(0.8),
                                content_width, Inches(0.3),
                                font_size=11, bold=False,
                                color='#64748b', alignment=PP_ALIGN.LEFT)
        
        # Box's M Test explanation box
        boxes_m_y = Inches(1.2)
        boxes_m_text = ""
        bg_light = soup.find('div', class_='bg-light')
        if bg_light:
            boxes_m_text = clean_text(bg_light.get_text())
        else:
            # Look for Box's M Test text
            for p in soup.find_all('p'):
                text = clean_text(p.get_text())
                if 'Box\'s M Test' in text or 'covariance matrices' in text.lower():
                    boxes_m_text = text
                    break
        
        if boxes_m_text:
            add_text_element(slide, boxes_m_text,
                            left_margin + Inches(0.1), boxes_m_y,
                            content_width - Inches(0.2), Inches(0.5),
                            font_size=8, bold=False,
                            color='#334155', alignment=PP_ALIGN.LEFT)
        
        # Two column layout for cards
        left_card_left = left_margin
        left_card_width = Inches(4.5)
        right_card_left = Inches(5.5)
        right_card_width = Inches(4)
        card_y = Inches(1.8)
        card_height = Inches(3.5)
        
        # Find the two cards
        cards = soup.find_all(class_='card')
        if not cards:
            # Look in two-col div
            two_col = soup.find(class_='two-col')
            if two_col:
                cards = two_col.find_all('div', class_='card', recursive=False)
        
        # Left card: "What Theory Said"
        theory_card = None
        if cards and len(cards) >= 1:
            theory_card = cards[0]
        
        theory_header = "What Theory Said"
        theory_content = []
        
        if theory_card:
            h3 = theory_card.find('h3')
            if h3:
                theory_header = clean_text(h3.get_text())
            
            # Extract all paragraphs and list items
            for p in theory_card.find_all('p'):
                text = clean_text(p.get_text())
                if text and len(text) > 10:
                    theory_content.append(text)
            for li in theory_card.find_all('li'):
                text = clean_text(li.get_text())
                if text and len(text) > 10:
                    theory_content.append(text)
        
        # Default theory content
        if not theory_content:
            theory_content = [
                "Box's M Test: p < 0.05",
                "We rejected the null hypothesis of equal covariance. This means the covariance matrices differ significantly between diabetes and no-diabetes groups.",
                "Statistical theory clearly suggests QDA should win (since it allows different covariances).",
                "The textbooks were confident..."
            ]
        
        # Right card: "What Data Revealed"
        reality_card = None
        if cards and len(cards) >= 2:
            reality_card = cards[1]
        
        reality_header = "What Data Revealed"
        reality_content = []
        
        if reality_card:
            h3 = reality_card.find('h3')
            if h3:
                reality_header = clean_text(h3.get_text())
            
            for p in reality_card.find_all('p'):
                text = clean_text(p.get_text())
                if text and len(text) > 10:
                    reality_content.append(text)
            for li in reality_card.find_all('li'):
                text = clean_text(li.get_text())
                if text and len(text) > 10:
                    reality_content.append(text)
        
        # Default reality content
        if not reality_content:
            reality_content = [
                "LDA Won by 4.79%",
                "Despite the violation, LDA was more robust. QDA likely overfitted the noise in the smaller class.",
                "...but reality had other plans"
            ]
        
        # Draw left card (Theory - red)
        theory_y = card_y
        add_text_element(slide, theory_header,
                        left_card_left + Inches(0.1), theory_y,
                        left_card_width - Inches(0.2), Inches(0.35),
                        font_size=14, bold=True,
                        color='#ef4444', alignment=PP_ALIGN.LEFT)
        theory_y += Inches(0.4)
        
        for item in theory_content[:4]:
            if theory_y < Inches(5.5):
                add_text_element(slide, f"• {item}" if not item.startswith('•') else item,
                                left_card_left + Inches(0.15), theory_y,
                                left_card_width - Inches(0.3), Inches(0.4),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                theory_y += Inches(0.45)
        
        # Draw right card (Reality - green)
        reality_y = card_y
        add_text_element(slide, reality_header,
                        right_card_left + Inches(0.1), reality_y,
                        right_card_width - Inches(0.2), Inches(0.35),
                        font_size=14, bold=True,
                        color='#10b981', alignment=PP_ALIGN.LEFT)
        reality_y += Inches(0.4)
        
        for item in reality_content[:3]:
            if reality_y < Inches(5.5):
                add_text_element(slide, f"• {item}" if not item.startswith('•') else item,
                                right_card_left + Inches(0.15), reality_y,
                                right_card_width - Inches(0.3), Inches(0.4),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                reality_y += Inches(0.45)
        
        # Key Insight box at bottom (if exists)
        key_insight_y = Inches(5.6)
        key_insight_text = ""
        
        # Look for Key Insight
        for div in soup.find_all('div'):
            text = clean_text(div.get_text())
            if 'Key Insight' in text:
                key_insight_text = text
                break
        
        if key_insight_text:
            add_text_element(slide, key_insight_text,
                            left_margin + Inches(0.1), key_insight_y,
                            content_width - Inches(0.2), Inches(0.6),
                            font_size=9, bold=True,
                            color='#005088', alignment=PP_ALIGN.LEFT)
        
        return
    
    # Special handling for "Interactive Dashboard: See It Live" slide - image left, feature boxes right
    if is_dashboard_slide:
        # Title at top
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.2), 
                            content_width, Inches(0.5),
                            font_size=24, bold=True, 
                            color='#005088', alignment=PP_ALIGN.LEFT)
        
        # Subtitle
        subtitle = soup.find('p')
        if subtitle:
            subtitle_text = clean_text(subtitle.get_text())
            if 'class balance' in subtitle_text.lower() or 'smote' in subtitle_text.lower():
                add_text_element(slide, subtitle_text,
                                left_margin, Inches(0.7),
                                content_width, Inches(0.25),
                                font_size=10, bold=False,
                                color='#64748b', alignment=PP_ALIGN.LEFT)
        
        # Two column layout
        left_col_width = Inches(5)
        right_col_left = Inches(5.5)
        right_col_width = Inches(4)
        col_y = Inches(1.1)
        
        # Left column: Dashboard image
        if slide_data['images']:
            img_data = slide_data['images'][0]
            img_path = img_data['path']
            if os.path.exists(img_path):
                try:
                    width, height = get_image_size(img_path, left_col_width, Inches(4.0))
                    slide.shapes.add_picture(img_path, left_margin, col_y, width, height)
                    
                    # Caption below image
                    caption = soup.find(class_='image-caption')
                    if caption:
                        caption_text = clean_text(caption.get_text())
                        add_text_element(slide, caption_text,
                                        left_margin, col_y + height + Inches(0.05),
                                        left_col_width, Inches(0.25),
                                        font_size=7, bold=False,
                                        color='#64748b', alignment=PP_ALIGN.CENTER)
                except Exception as e:
                    print(f"    Warning: Could not add dashboard image {img_path}: {e}")
        
        # Right column: Two feature boxes
        box_y = col_y
        box_height = Inches(1.8)
        box_gap = Inches(0.15)
        
        # Find the two feature cards
        cards = soup.find_all(class_='card')
        if not cards:
            # Look in the right column div
            two_col = soup.find(class_='two-col')
            if two_col:
                right_col_div = two_col.find_all('div', recursive=False)
                if len(right_col_div) >= 2:
                    cards = right_col_div[1].find_all('div', class_='card', recursive=False)
        
        # Box 1: Interactive Features
        features_card = None
        if cards and len(cards) >= 1:
            features_card = cards[0]
        
        features_header = "Interactive Features"
        features_items = []
        
        if features_card:
            h4 = features_card.find('h4')
            if h4:
                features_header = clean_text(h4.get_text())
                # Remove icon text if present
                features_header = features_header.replace('🛠️', '').replace('⚙️', '').strip()
            
            ul = features_card.find('ul')
            if ul:
                for li in ul.find_all('li'):
                    text = clean_text(li.get_text())
                    if text:
                        features_items.append(text)
        
        # Default features
        if not features_items:
            features_items = [
                "Class Distribution Slider: Adjust from 40% to 70% to see how balance affects performance",
                "SMOTE Toggle: Enable/disable synthetic oversampling",
                "Real-Time Updates: Metrics recalculate with actual R model training",
                "Live ROC Curves: Watch curves update as parameters change"
            ]
        
        # Box 2: What You Can Explore
        explore_card = None
        if cards and len(cards) >= 2:
            explore_card = cards[1]
        
        explore_header = "What You Can Explore"
        explore_items = []
        
        if explore_card:
            h4 = explore_card.find('h4')
            if h4:
                explore_header = clean_text(h4.get_text())
                explore_header = explore_header.replace('🔍', '').replace('📊', '').strip()
            
            ul = explore_card.find('ul')
            if ul:
                for li in ul.find_all('li'):
                    text = clean_text(li.get_text())
                    if text:
                        explore_items.append(text)
        
        # Default explore items
        if not explore_items:
            explore_items = [
                "Balanced Classes (50/50): See how specificity improves",
                "SMOTE Impact: Compare with/without synthetic oversampling",
                "Model Comparison: Watch LDA vs QDA performance change",
                "Tooltips: Hover over metrics for explanations"
            ]
        
        # Draw Box 1: Interactive Features
        add_text_element(slide, features_header,
                        right_col_left + Inches(0.1), box_y,
                        right_col_width - Inches(0.2), Inches(0.3),
                        font_size=12, bold=True,
                        color='#005088', alignment=PP_ALIGN.LEFT)
        box_y += Inches(0.35)
        
        for item in features_items[:4]:
            if box_y < Inches(3.0):
                add_text_element(slide, f"• {item}" if not item.startswith('•') else item,
                                right_col_left + Inches(0.15), box_y,
                                right_col_width - Inches(0.3), Inches(0.3),
                                font_size=8, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                box_y += Inches(0.32)
        
        box_y += Inches(0.1)
        
        # Draw Box 2: What You Can Explore
        add_text_element(slide, explore_header,
                        right_col_left + Inches(0.1), box_y,
                        right_col_width - Inches(0.2), Inches(0.3),
                        font_size=12, bold=True,
                        color='#005088', alignment=PP_ALIGN.LEFT)
        box_y += Inches(0.35)
        
        for item in explore_items[:4]:
            if box_y < Inches(5.5):
                add_text_element(slide, f"• {item}" if not item.startswith('•') else item,
                                right_col_left + Inches(0.15), box_y,
                                right_col_width - Inches(0.3), Inches(0.3),
                                font_size=8, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                box_y += Inches(0.32)
        
        # Live Demo section at bottom
        demo_y = Inches(5.6)
        
        # Find Live Demo link
        demo_text = ""
        demo_url = ""
        for a in soup.find_all('a'):
            href = a.get('href', '')
            if 'vercel.app' in href or 'dashboard' in href.lower():
                demo_url = href
                demo_text = clean_text(a.get_text())
                break
        
        if not demo_url:
            demo_url = "https://asds-5303-final-project-presentatio.vercel.app/dashboard.html"
            demo_text = demo_url
        
        # Live Demo header
        add_text_element(slide, "Live Demo: Try the interactive dashboard yourself!",
                        left_margin + Inches(0.1), demo_y,
                        content_width - Inches(0.2), Inches(0.25),
                        font_size=10, bold=True,
                        color='#005088', alignment=PP_ALIGN.LEFT)
        demo_y += Inches(0.3)
        
        # URL
        add_text_element(slide, demo_url,
                        left_margin + Inches(0.1), demo_y,
                        content_width - Inches(0.2), Inches(0.25),
                        font_size=9, bold=False,
                        color='#2563eb', alignment=PP_ALIGN.LEFT)
        demo_y += Inches(0.3)
        
        # Instruction text
        instruction = "Adjust the slider and toggle SMOTE to see real-time model performance changes."
        add_text_element(slide, instruction,
                        left_margin + Inches(0.1), demo_y,
                        content_width - Inches(0.2), Inches(0.25),
                        font_size=8, bold=False,
                        color='#64748b', alignment=PP_ALIGN.LEFT)
        
        return
    
    # Special handling for "What We Learned: The Journey's End" slide - two columns and summary boxes
    if is_learned_slide:
        # Title at top
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.2), 
                            content_width, Inches(0.6),
                            font_size=28, bold=True, 
                            color='#005088', alignment=PP_ALIGN.LEFT)
        
        # Two column layout
        left_col_width = Inches(4.5)
        right_col_left = Inches(5.5)
        right_col_width = Inches(4)
        col_y = Inches(1.0)
        
        # Find the two columns
        two_col = soup.find(class_='two-col')
        cards = []
        if two_col:
            cards = two_col.find_all('div', class_='card', recursive=False)
        
        # Left column: Key Takeaways
        takeaways_card = None
        if cards and len(cards) >= 1:
            takeaways_card = cards[0]
        
        takeaways_header = "Key Takeaways"
        takeaways_items = []
        
        if takeaways_card:
            h4 = takeaways_card.find('h4')
            if h4:
                takeaways_header = clean_text(h4.get_text())
                # Remove icon text if present
                takeaways_header = takeaways_header.replace('🏆', '').replace('📊', '').strip()
            
            ul = takeaways_card.find('ul')
            if ul:
                for li in ul.find_all('li'):
                    text = clean_text(li.get_text())
                    if text:
                        takeaways_items.append(text)
            else:
                # Try paragraphs
                for p in takeaways_card.find_all('p'):
                    text = clean_text(p.get_text())
                    if text and len(text) > 20:
                        takeaways_items.append(text)
        
        # Default takeaways
        if not takeaways_items:
            takeaways_items = [
                "LDA Wins: Despite theory favoring QDA, LDA achieved 76.09% accuracy with superior robustness. Sometimes simpler is better.",
                "Critical Gap: Missing 44% of diabetic cases is unacceptable. We must do better.",
                "Root Cause: Class imbalance (65% vs 35%) is the villain in our story. This must be addressed."
            ]
        
        # Right column: The Path Forward
        pathforward_card = None
        if cards and len(cards) >= 2:
            pathforward_card = cards[1]
        
        pathforward_header = "The Path Forward"
        pathforward_items = []
        
        if pathforward_card:
            h4 = pathforward_card.find('h4')
            if h4:
                pathforward_header = clean_text(h4.get_text())
                pathforward_header = pathforward_header.replace('🚀', '').replace('🔮', '').strip()
            
            ul = pathforward_card.find('ul')
            if ul:
                for li in ul.find_all('li'):
                    text = clean_text(li.get_text())
                    if text:
                        pathforward_items.append(text)
            else:
                # Try paragraphs
                for p in pathforward_card.find_all('p'):
                    text = clean_text(p.get_text())
                    if text and len(text) > 20:
                        pathforward_items.append(text)
        
        # Default path forward items
        if not pathforward_items:
            pathforward_items = [
                "SMOTE: Balance the classes to give diabetes cases a fair chance.",
                "Cost-Sensitive Learning: Prioritize catching diabetic cases—missing them costs more than false alarms.",
                "Ensemble Methods: Combine LDA with Random Forest, SVM for a stronger team.",
                "Feature Engineering: Uncover hidden patterns through interaction terms."
            ]
        
        # Draw left column: Key Takeaways
        takeaways_y = col_y
        add_text_element(slide, takeaways_header,
                        left_margin + Inches(0.1), takeaways_y,
                        left_col_width - Inches(0.2), Inches(0.35),
                        font_size=14, bold=True,
                        color='#005088', alignment=PP_ALIGN.LEFT)
        takeaways_y += Inches(0.4)
        
        for item in takeaways_items[:3]:
            if takeaways_y < Inches(4.5):
                add_text_element(slide, f"• {item}" if not item.startswith('•') else item,
                                left_margin + Inches(0.15), takeaways_y,
                                left_col_width - Inches(0.3), Inches(0.5),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                takeaways_y += Inches(0.6)
        
        # Draw right column: The Path Forward
        pathforward_y = col_y
        add_text_element(slide, pathforward_header,
                        right_col_left + Inches(0.1), pathforward_y,
                        right_col_width - Inches(0.2), Inches(0.35),
                        font_size=14, bold=True,
                        color='#005088', alignment=PP_ALIGN.LEFT)
        pathforward_y += Inches(0.4)
        
        for item in pathforward_items[:4]:
            if pathforward_y < Inches(4.5):
                add_text_element(slide, f"• {item}" if not item.startswith('•') else item,
                                right_col_left + Inches(0.15), pathforward_y,
                                right_col_width - Inches(0.3), Inches(0.4),
                                font_size=9, bold=False,
                                color='#334155', alignment=PP_ALIGN.LEFT)
                pathforward_y += Inches(0.5)
        
        # Bottom summary boxes
        summary_y = Inches(5.2)
        
        # Find "The Bottom Line" and "Our Promise" boxes
        bottom_line_text = ""
        our_promise_text = ""
        
        # Look for boxes with specific background colors or text
        for div in soup.find_all('div'):
            text = clean_text(div.get_text())
            style = div.get('style', '')
            if 'Bottom Line' in text or 'not ready for clinical' in text.lower():
                bottom_line_text = text
            elif 'Our Promise' in text or 'continue improving' in text.lower():
                our_promise_text = text
        
        # Default texts
        if not bottom_line_text:
            bottom_line_text = "The Bottom Line: We have a good start, but we're not ready for clinical deployment yet."
        if not our_promise_text:
            our_promise_text = "Our Promise: We'll continue improving until we can confidently detect diabetes early."
        
        # Draw "The Bottom Line" box (light blue)
        add_text_element(slide, bottom_line_text,
                        left_margin + Inches(0.1), summary_y,
                        content_width - Inches(0.2), Inches(0.4),
                        font_size=10, bold=True,
                        color='#005088', alignment=PP_ALIGN.LEFT)
        summary_y += Inches(0.45)
        
        # Draw "Our Promise" box (light green)
        add_text_element(slide, our_promise_text,
                        left_margin + Inches(0.1), summary_y,
                        content_width - Inches(0.2), Inches(0.4),
                        font_size=10, bold=True,
                        color='#10b981', alignment=PP_ALIGN.LEFT)
        
        return
    
    # Special handling for "Questions?" slide with References - dark blue background
    if is_questions_slide:
        # Set background color to dark blue
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 80, 136)  # #005088
        
        # Title at top center (white)
        if slide_data['title']:
            add_text_element(slide, slide_data['title'], 
                            left_margin, Inches(0.5), 
                            content_width, Inches(1.0),
                            font_size=48, bold=True, 
                            color='#ffffff', alignment=PP_ALIGN.CENTER)
        
        # References box
        ref_y = Inches(2.0)
        ref_width = Inches(8.5)
        ref_left = left_margin + (content_width - ref_width) / 2
        
        # Find References section
        references_text = []
        ref_header = "References:"
        
        # Look for References text
        for div in soup.find_all('div'):
            text = clean_text(div.get_text())
            if 'References' in text or 'references' in text.lower():
                # Extract all paragraphs or list items
                for p in div.find_all('p'):
                    p_text = clean_text(p.get_text())
                    if p_text and len(p_text) > 20 and 'References' not in p_text:
                        references_text.append(p_text)
                for li in div.find_all('li'):
                    li_text = clean_text(li.get_text())
                    if li_text and len(li_text) > 20:
                        references_text.append(li_text)
                # Also try to get all text and split by numbers
                full_text = text
                if 'References' in full_text:
                    # Split by numbered items
                    import re
                    items = re.split(r'\n\s*(\d+\.)', full_text)
                    for i in range(1, len(items), 2):
                        if i + 1 < len(items):
                            ref_item = items[i] + items[i+1]
                            ref_item = clean_text(ref_item)
                            if ref_item and len(ref_item) > 20:
                                references_text.append(ref_item)
                break
        
        # Default references if not found
        if not references_text:
            references_text = [
                "1. International Diabetes Federation. (2021). IDF Diabetes Atlas, 10th Edition. Brussels, Belgium: International Diabetes Federation.",
                "2. National Institute of Diabetes and Digestive and Kidney Diseases. Pima Indians Diabetes Database. UCI Machine Learning Repository.",
                "3. James, G., Witten, D., Hastie, T., & Tibshirani, R. (2013). An Introduction to Statistical Learning: with Applications in R. New York: Springer.",
                "4. Venables, W. N., & Ripley, B. D. (2002). Modern Applied Statistics with S (4th ed.). New York: Springer.",
                "5. American Diabetes Association. (2023). Standards of Care in Diabetes – 2023. Diabetes Care, 46(Supplement 1), S1-S291.",
                "6. Chawla, N. V., Bowyer, K. W., Hall, L. O., & Kegelmeyer, W. P. (2002). SMOTE: Synthetic Minority Over-sampling Technique. Journal of Artificial Intelligence Research, 16, 321-357."
            ]
        
        # Draw References header
        add_text_element(slide, ref_header,
                        ref_left + Inches(0.1), ref_y,
                        ref_width - Inches(0.2), Inches(0.35),
                        font_size=14, bold=True,
                        color='#ffffff', alignment=PP_ALIGN.LEFT)
        ref_y += Inches(0.4)
        
        # Draw references list
        for ref in references_text[:6]:
            if ref_y < Inches(6.5):
                # Ensure it starts with a number if it doesn't
                if not ref[0].isdigit():
                    # Try to extract number from text
                    import re
                    match = re.search(r'(\d+)\.', ref)
                    if not match:
                        ref = f"{len(references_text[:6]) - references_text.index(ref) + 1}. {ref}"
                
                add_text_element(slide, ref,
                                ref_left + Inches(0.1), ref_y,
                                ref_width - Inches(0.2), Inches(0.4),
                                font_size=9, bold=False,
                                color='#e0f2fe', alignment=PP_ALIGN.LEFT)
                ref_y += Inches(0.45)
        
        return
    
    # Add title if exists
    if slide_data['title']:
        add_text_element(slide, slide_data['title'], 
                        left_margin, Inches(0.2), 
                        content_width, Inches(0.6),
                        font_size=28, bold=True, 
                        color='#005088', alignment=PP_ALIGN.LEFT)
        y_pos = Inches(1.0)
    
    # Special handling for title slide
    if is_title_slide:
        # Special handling for title slide
        y_pos = Inches(2.5)
        # Title is already added, now add subtitle and authors
        paragraphs = soup.find_all(['p'])
        for para in paragraphs[:5]:
            text = clean_text(para.get_text())
            if text and y_pos < max_y:
                add_text_element(slide, text,
                                left_margin, y_pos,
                                content_width, Inches(0.3),
                                font_size=14, alignment=PP_ALIGN.LEFT)
                y_pos += Inches(0.4)
        return
    
    # Handle images first (but not for hook slide which is text-only)
    if slide_data['images'] and not is_hook_slide:
        for img_data in slide_data['images']:
            if y_pos > max_y - Inches(3):
                break
            
            img_path = img_data['path']
            if os.path.exists(img_path):
                try:
                    width, height = get_image_size(img_path, Inches(7), Inches(4))
                    # Center image horizontally
                    img_left = (Inches(10) - width) / 2
                    slide.shapes.add_picture(img_path, img_left, y_pos, width, height)
                    y_pos += height + Inches(0.2)
                except Exception as e:
                    print(f"    Warning: Could not add image {img_path}: {e}")
    
    # Process paragraphs with better formatting detection
    paragraphs = soup.find_all(['p'])
    for elem in paragraphs[:15]:  # Limit to avoid overflow
        text = clean_text(elem.get_text())
        if not text or text == slide_data['title'] or len(text) < 3:
            continue
        
        # Skip if it's an image caption (usually very short)
        if len(text) < 50 and ('png' in text.lower() or 'jpg' in text.lower()):
            continue
        
        # Check if we have space
        if y_pos > max_y:
            break
        
        # Detect formatting from inline styles or classes
        font_size = 11
        bold = False
        color = '#334155'
        alignment = PP_ALIGN.LEFT
        
        # Check for style attributes
        style = elem.get('style', '')
        if 'font-size' in style:
            # Extract font size (simplified)
            if '0.9em' in style or '0.85em' in style:
                font_size = 10
            elif '0.75em' in style or '0.7em' in style:
                font_size = 9
            elif '0.65em' in style or '0.6em' in style:
                font_size = 8
        
        # Check for color in style
        if 'color:' in style:
            if '#005088' in style or '#005' in style:
                color = '#005088'
            elif '#ef4444' in style or '#ef4' in style:
                color = '#ef4444'
            elif '#64748b' in style:
                color = '#64748b'
        
        # Check for bold
        if elem.find('strong') or '<strong>' in str(elem):
            bold = True
        
        # Check for center alignment
        if 'text-align: center' in style or 'text-align:center' in style:
            alignment = PP_ALIGN.CENTER
            left_margin = Inches(1)
            content_width = Inches(8)
        
        height = Inches(0.3)
        add_text_element(slide, text,
                        left_margin, y_pos,
                        content_width, height,
                        font_size=font_size, bold=bold,
                        color=color, alignment=alignment)
        y_pos += height + Inches(0.1)
        
        # Reset margins for next element
        left_margin = Inches(0.5)
        content_width = Inches(9)
    
    # Process headings (h3, h4)
    headings = soup.find_all(['h3', 'h4'])
    for elem in headings[:5]:
        text = clean_text(elem.get_text())
        if not text or text == slide_data['title'] or y_pos > max_y:
            continue
        
        add_text_element(slide, text,
                        left_margin, y_pos,
                        content_width, Inches(0.35),
                        font_size=16, bold=True,
                        color='#005088', alignment=PP_ALIGN.LEFT)
        y_pos += Inches(0.4)
    
    # Process lists
    lists = soup.find_all(['ul', 'ol'])
    for ul in lists[:3]:  # Limit lists
        items = ul.find_all('li', recursive=False)
        for item in items[:6]:  # Limit items per list
            text = clean_text(item.get_text())
            if not text or y_pos > max_y:
                break
            
            # Add bullet point
            add_text_element(slide, f"• {text}",
                            left_margin + Inches(0.2), y_pos,
                            content_width - Inches(0.2), Inches(0.25),
                            font_size=10, alignment=PP_ALIGN.LEFT)
            y_pos += Inches(0.3)
    
    # Process cards and special sections
    cards = soup.find_all(class_=['card', 'bg-light', 'callout-box', 'story-highlight'])
    for card in cards[:2]:  # Limit cards
        text = clean_text(card.get_text())
        if text and y_pos < max_y and len(text) > 10:
            add_text_element(slide, text,
                            left_margin, y_pos,
                            content_width, Inches(0.5),
                            font_size=10, alignment=PP_ALIGN.LEFT)
            y_pos += Inches(0.6)

def create_pptx_from_slides(slides, output_file):
    """Create PowerPoint presentation from slides"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print(f"Creating {len(slides)} slides...")
    
    for i, slide_data in enumerate(slides, 1):
        title_preview = slide_data['title'][:50] if slide_data['title'] else "Untitled"
        print(f"  Processing slide {i}/{len(slides)}: {title_preview}...")
        if slide_data['images']:
            print(f"    Found {len(slide_data['images'])} image(s)")
        
        # Create blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Set background color if specified
        if slide_data['background_color']:
            try:
                background = slide.background
                fill = background.fill
                fill.solid()
                r, g, b = hex_to_rgb(slide_data['background_color'])
                fill.fore_color.rgb = RGBColor(r, g, b)
            except Exception as e:
                print(f"    Warning: Could not set background color: {e}")
        
        # Process and add content
        try:
            process_slide_content(slide, slide_data)
        except Exception as e:
            print(f"    Warning: Error processing slide content: {e}")
    
    prs.save(output_file)
    print(f"\n✅ PowerPoint presentation saved to: {output_file}")

def main():
    html_file = 'final_presentation_ASDS_5303.html'
    output_file = 'ASDS_5303_Final_Presentation_Group_7.pptx'
    images_dir = 'images'
    
    if not os.path.exists(html_file):
        print(f"❌ Error: {html_file} not found!")
        return
    
    if not os.path.exists(images_dir):
        print(f"⚠️  Warning: {images_dir} directory not found. Images will be skipped.")
    
    print(f"📖 Parsing {html_file}...")
    try:
        slides = parse_html_slides(html_file)
        print(f"📊 Found {len(slides)} slides")
        
        # Count total images
        total_images = sum(len(s['images']) for s in slides)
        print(f"🖼️  Found {total_images} image references")
        
        print(f"📝 Creating PowerPoint presentation...")
        create_pptx_from_slides(slides, output_file)
        
        print(f"\n✅ Conversion complete!")
        print(f"📄 Output file: {output_file}")
        print(f"\n💡 Note: Review the presentation and adjust formatting as needed.")
    except Exception as e:
        print(f"❌ Error during conversion: {e}")
        import traceback
        traceback.print_exc()

if __name__ == '__main__':
    main()
